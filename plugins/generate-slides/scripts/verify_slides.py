#!/usr/bin/env python3
"""Verification harness for generated PPTX slide decks.

Performs deterministic quality checks on a .pptx file and reports
pass/fail for each criterion.  Designed to run as a post-generation
gate in the generate-slides skill.

Usage:
    python verify_slides.py output.pptx [--style slide-style.yaml]

Exit code 0 = all checks pass, 1 = one or more failures.
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# ── Helpers ──────────────────────────────────────────────────────────

SLIDE_W_EMU = 12192000   # 13.333" in EMU
SLIDE_H_EMU = 6858000    #  7.5"   in EMU
MARGIN_EMU = 731520      #  0.8"   in EMU (left/right safe zone)
TOP_MARGIN_EMU = 137160  #  0.15"  in EMU (title top)
BOT_MARGIN_EMU = 274320  #  0.3"   in EMU (bottom safe zone)
OVERLAP_TOLERANCE_EMU = 45720  # 0.05" — shapes closer than this are "overlapping"

# Words that strongly suggest a topic-label title (not a sentence).
TOPIC_LABEL_PATTERNS = [
    r"^results?$", r"^methods?$", r"^background$", r"^introduction$",
    r"^discussion$", r"^overview$", r"^summary$", r"^conclusions?$",
    r"^motivation$", r"^outline$", r"^agenda$", r"^objectives?$",
    r"^approach$", r"^next steps$", r"^future work$",
    r"^data$", r"^analysis$", r"^evaluation$", r"^implementation$",
    r"^model performance$", r"^experimental setup$",
]

THANK_YOU_PATTERNS = [
    r"thank\s*you", r"thanks?[\!\.]?$", r"questions\??$",
    r"q\s*&\s*a\??$", r"any questions\??$",
]


def _get_slide_title(slide):
    """Return the first text box's text as the slide title, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return None


def _is_topic_label(title):
    """Check if a title looks like a topic label instead of a sentence."""
    t = title.strip().rstrip(".!?").lower()
    # Very short titles (1-3 words) without a verb are likely labels
    for pat in TOPIC_LABEL_PATTERNS:
        if re.match(pat, t, re.IGNORECASE):
            return True
    return False


def _is_sentence(title):
    """Heuristic: a sentence has >3 words or contains a verb-like structure."""
    words = title.strip().split()
    if len(words) >= 5:
        return True
    # 3-4 word titles can be sentences if they have a verb
    # but we can't reliably detect verbs without NLP, so we're lenient
    return len(words) >= 4


# ── Checks ───────────────────────────────────────────────────────────

class CheckResult:
    def __init__(self, name, passed, details=""):
        self.name = name
        self.passed = passed
        self.details = details

    def __str__(self):
        status = "PASS" if self.passed else "FAIL"
        s = f"  [{status}] {self.name}"
        if self.details:
            s += f"\n         {self.details}"
        return s


def check_action_titles(prs):
    """Verify slide titles are action titles (sentences), not topic labels."""
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        title = _get_slide_title(slide)
        if title is None:
            continue
        # Skip title slide (slide 1) — it's the presentation title, not an action title
        if i == 1:
            continue
        if _is_topic_label(title):
            issues.append(f"Slide {i}: \"{title}\" looks like a topic label, not an action title")
    if issues:
        return CheckResult("Action titles", False, "; ".join(issues))
    return CheckResult("Action titles", True)


def check_ghost_deck(prs):
    """Print all titles sequentially for the ghost deck test."""
    titles = []
    for i, slide in enumerate(prs.slides, 1):
        title = _get_slide_title(slide)
        if title:
            titles.append(f"  {i}. {title}")
    ghost = "\n".join(titles)
    # This is informational — always passes but prints the ghost deck
    return CheckResult(
        "Ghost deck test (informational)",
        True,
        f"Read these titles in sequence — do they tell a complete story?\n{ghost}"
    )


def check_figure_bounds(prs):
    """Verify no images overflow the slide edges."""
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                right = shape.left + shape.width
                bottom = shape.top + shape.height
                if right > SLIDE_W_EMU:
                    overflow_in = (right - SLIDE_W_EMU) / 914400
                    issues.append(f"Slide {i}: image overflows right edge by {overflow_in:.2f}\"")
                if bottom > SLIDE_H_EMU:
                    overflow_in = (bottom - SLIDE_H_EMU) / 914400
                    issues.append(f"Slide {i}: image overflows bottom edge by {overflow_in:.2f}\"")
    if issues:
        return CheckResult("Figure bounds", False, "; ".join(issues))
    return CheckResult("Figure bounds", True)


def check_font_consistency(prs, expected_font=None):
    """Check that all text uses the expected font family."""
    if expected_font is None:
        expected_font = "Avenir Light"
    off_font = {}
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font.name
                    if font and font != expected_font and font != "PT Mono":
                        key = f"Slide {i}: \"{font}\""
                        off_font[key] = off_font.get(key, 0) + 1
    if off_font:
        details = "; ".join(f"{k} ({v} runs)" for k, v in off_font.items())
        return CheckResult("Font consistency", False, f"Expected \"{expected_font}\", found: {details}")
    return CheckResult("Font consistency", True)


def check_no_thank_you_ending(prs):
    """Verify the last main slide is not 'Thank You' or 'Questions?'."""
    if len(prs.slides) == 0:
        return CheckResult("Conclusions ending", True, "No slides")
    last_title = _get_slide_title(prs.slides[-1])
    if last_title:
        for pat in THANK_YOU_PATTERNS:
            if re.search(pat, last_title, re.IGNORECASE):
                return CheckResult(
                    "Conclusions ending", False,
                    f"Last slide title is \"{last_title}\" — should end with Conclusions, not Thank You/Questions"
                )
    return CheckResult("Conclusions ending", True)


def check_content_density(prs, max_bullets=6):
    """Flag slides with too many bullets (>max_bullets)."""
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bullet_count = sum(
                1 for p in shape.text_frame.paragraphs
                if p.text.strip().startswith("\u2022") or p.level > 0
            )
            if bullet_count > max_bullets:
                issues.append(f"Slide {i}: {bullet_count} bullets (max {max_bullets})")
    if issues:
        return CheckResult("Content density", False, "; ".join(issues))
    return CheckResult("Content density", True)


def check_shape_overlaps(prs):
    """Detect shapes that overlap each other (e.g., title vs accent line, text vs figure).

    Two shapes overlap if their bounding boxes intersect with more than
    OVERLAP_TOLERANCE_EMU of penetration in both axes.
    """
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        for a_idx in range(len(shapes)):
            for b_idx in range(a_idx + 1, len(shapes)):
                a, b = shapes[a_idx], shapes[b_idx]
                # Compute overlap in x and y
                x_overlap = min(a.left + a.width, b.left + b.width) - max(a.left, b.left)
                y_overlap = min(a.top + a.height, b.top + b.height) - max(a.top, b.top)
                if x_overlap > OVERLAP_TOLERANCE_EMU and y_overlap > OVERLAP_TOLERANCE_EMU:
                    # Identify shape types for the message
                    def _desc(s):
                        if s.shape_type == 13:
                            return "image"
                        if hasattr(s, "text") and s.text:
                            text_preview = s.text[:30].replace("\n", " ")
                            return f"text(\"{text_preview}...\")" if len(s.text) > 30 else f"text(\"{s.text}\")"
                        return f"shape({s.shape_type})"
                    overlap_in = min(x_overlap, y_overlap) / 914400
                    issues.append(
                        f"Slide {i}: {_desc(a)} overlaps {_desc(b)} by {overlap_in:.2f}\""
                    )
    if issues:
        return CheckResult("Shape overlaps", False, "; ".join(issues[:5]))  # cap at 5
    return CheckResult("Shape overlaps", True)


def check_margin_compliance(prs):
    """Verify all shapes stay within the slide margins (not just slide edges).

    Content should respect:
    - Left margin: 0.8" (MARGIN_EMU)
    - Right margin: slide_width - 0.8"
    - Bottom margin: slide_height - 0.3"

    The title slide (slide 1) uses different positioning and is skipped.
    """
    issues = []
    right_limit = SLIDE_W_EMU - MARGIN_EMU
    bottom_limit = SLIDE_H_EMU - BOT_MARGIN_EMU

    for i, slide in enumerate(prs.slides, 1):
        if i == 1:  # title slide has its own layout
            continue
        for shape in slide.shapes:
            shape_right = shape.left + shape.width
            shape_bottom = shape.top + shape.height

            if shape.left < MARGIN_EMU - OVERLAP_TOLERANCE_EMU:
                intrusion = (MARGIN_EMU - shape.left) / 914400
                desc = "image" if shape.shape_type == 13 else "shape"
                issues.append(f"Slide {i}: {desc} starts {intrusion:.2f}\" inside left margin")

            if shape_right > right_limit + OVERLAP_TOLERANCE_EMU:
                intrusion = (shape_right - right_limit) / 914400
                desc = "image" if shape.shape_type == 13 else "shape"
                issues.append(f"Slide {i}: {desc} extends {intrusion:.2f}\" past right margin")

            if shape_bottom > bottom_limit + OVERLAP_TOLERANCE_EMU:
                intrusion = (shape_bottom - bottom_limit) / 914400
                desc = "image" if shape.shape_type == 13 else "shape"
                issues.append(f"Slide {i}: {desc} extends {intrusion:.2f}\" past bottom margin")

    if issues:
        return CheckResult("Margin compliance", False, "; ".join(issues[:5]))
    return CheckResult("Margin compliance", True)


def check_text_alignment(prs):
    """Flag text boxes with inconsistent left-edge alignment on the same slide.

    Good slides have text boxes snapped to a grid (typically LEFT_MARGIN).
    If text boxes on the same slide have left edges that differ by a small
    but noticeable amount (0.1"–0.5"), they look sloppy.
    """
    issues = []
    alignment_tolerance = 91440  # 0.1" — tighter than overlap tolerance

    for i, slide in enumerate(prs.slides, 1):
        if i == 1:  # title slide
            continue
        # Collect left edges of all text-bearing shapes
        left_edges = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                left_edges.append(shape.left)

        if len(left_edges) < 2:
            continue

        # Group left edges by proximity
        left_edges.sort()
        groups = [[left_edges[0]]]
        for edge in left_edges[1:]:
            if edge - groups[-1][-1] <= alignment_tolerance:
                groups[-1].append(edge)
            else:
                groups.append([edge])

        # If there are >2 distinct alignment groups, some text is misaligned
        # (title + body = 2 groups is normal; 3+ suggests sloppy placement)
        if len(groups) > 2:
            edges_in = [f"{e / 914400:.2f}\"" for e in left_edges]
            issues.append(
                f"Slide {i}: {len(groups)} distinct left alignments ({', '.join(edges_in)}) — "
                f"text boxes may look misaligned"
            )

    if issues:
        return CheckResult("Text alignment", False, "; ".join(issues[:5]))
    return CheckResult("Text alignment", True)


def check_slide_count(prs):
    """Report slide count (informational)."""
    n = len(prs.slides)
    return CheckResult("Slide count", True, f"{n} slides")


# ── Main ─────────────────────────────────────────────────────────────

def verify(pptx_path, expected_font=None):
    """Run all checks and return (results, all_passed)."""
    prs = Presentation(pptx_path)
    results = [
        check_slide_count(prs),
        check_action_titles(prs),
        check_ghost_deck(prs),
        check_figure_bounds(prs),
        check_margin_compliance(prs),
        check_shape_overlaps(prs),
        check_text_alignment(prs),
        check_font_consistency(prs, expected_font),
        check_no_thank_you_ending(prs),
        check_content_density(prs),
    ]
    all_passed = all(r.passed for r in results)
    return results, all_passed


def main():
    parser = argparse.ArgumentParser(description="Verify a generated PPTX slide deck.")
    parser.add_argument("pptx", help="Path to the .pptx file to verify")
    parser.add_argument("--font", default=None, help="Expected body font (default: Avenir Light)")
    args = parser.parse_args()

    if not Path(args.pptx).exists():
        print(f"ERROR: File not found: {args.pptx}")
        sys.exit(1)

    print(f"Verifying: {args.pptx}")
    print("=" * 60)

    results, all_passed = verify(args.pptx, expected_font=args.font)

    for r in results:
        print(r)

    print("=" * 60)
    if all_passed:
        print("ALL CHECKS PASSED")
    else:
        n_fail = sum(1 for r in results if not r.passed)
        print(f"{n_fail} CHECK(S) FAILED — fix issues above before presenting")

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
