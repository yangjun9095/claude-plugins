#!/usr/bin/env python3
"""Generate sample-deck.pptx to showcase the default generate-slides style.

Run:  python generate_sample_deck.py
Out:  sample-deck.pptx  (in the same directory)
"""

import sys
from pathlib import Path

# Add the scripts/ dir so we can import the helper library
_scripts = Path(__file__).resolve().parent.parent / "scripts"
sys.path.insert(0, str(_scripts))

from pptx_helpers import *  # noqa: E402, F403


def slide_01_title(prs):
    """Title slide."""
    slide = add_blank_slide(prs)
    make_title_slide(
        slide,
        title="Sample Slide Deck — Default Style Preview",
        author="Your Name",
        affiliation="Organization / Lab / Team",
        date="April 2026",
    )


def slide_02_content(prs):
    """Content slide with bullet points."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Action titles state the takeaway as a sentence")
    add_accent_line(slide)
    add_bullets(slide, LEFT_MARGIN, TOP_MARGIN, CONTENT_W, Inches(4.5), [
        ("Action titles: ", "every slide title is a complete sentence, not a topic label"),
        ("Blue accent prefixes: ", "key terms are highlighted without using bold"),
        "Simple bullets work too — just pass a plain string",
        ("No bold anywhere: ", "hierarchy comes from font size and spacing only"),
        ("Customizable: ", "drop a slide-style.yaml in your project to override any setting"),
    ])


def slide_03_figure_placeholder(prs):
    """Slide demonstrating figure placement with a generated placeholder."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Figures are centered and auto-clamped to slide bounds")
    add_accent_line(slide)

    # Draw a placeholder rectangle where a figure would go
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    fig_left = Inches(2.0)
    fig_top = Inches(1.5)
    fig_w = Inches(9.0)
    fig_h = Inches(4.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, fig_left, fig_top, fig_w, fig_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    shape.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    shape.line.width = Pt(1.5)

    # Label inside the placeholder
    add_textbox(
        slide, fig_left + Inches(2.5), fig_top + Inches(1.8),
        Inches(4.0), Inches(1.0),
        "[Your figure here]",
        font_size=SUBTITLE_SIZE, color=GREY,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, LEFT_MARGIN, Inches(6.3), CONTENT_W, Inches(0.8),
        "Figures are auto-scaled to never overflow slide edges. Supports PNG, JPG, and SVG.",
        font_size=FOOTNOTE_SIZE, color=GREY,
    )


def slide_04_table(prs):
    """Table slide."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Tables use light grey headers with no bold in any cell")
    add_accent_line(slide)

    data = [
        ["Method", "Accuracy", "F1 Score", "Runtime"],
        ["Baseline", "72.3%", "0.68", "12 min"],
        ["Model A", "85.1%", "0.82", "18 min"],
        ["Model B", "91.7%", "0.90", "24 min"],
        ["Ours", "94.2%", "0.93", "15 min"],
    ]
    add_table(
        slide,
        left=Inches(1.5), top=Inches(1.5),
        width=Inches(10.0), height=Inches(3.5),
        rows=5, cols=4,
        data=data,
        highlight_cells={
            (4, 1): GREEN,  # best accuracy
            (4, 2): GREEN,  # best F1
        },
    )

    add_textbox(
        slide, LEFT_MARGIN, Inches(5.5), CONTENT_W, Inches(0.6),
        "Green highlights call out best-in-class values. Red is available for negative values.",
        font_size=FOOTNOTE_SIZE, color=GREY,
    )


def slide_05_code_and_callout(prs):
    """Code box + callout box."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Code boxes and callout boxes provide structured emphasis")
    add_accent_line(slide)

    add_code_box(
        slide,
        left=LEFT_MARGIN, top=Inches(1.4),
        width=Inches(5.5), height=Inches(2.8),
        text=(
            "# Code blocks use PT Mono on light grey\n"
            "from pptx_helpers import *\n"
            "\n"
            "prs = new_presentation()\n"
            "slide = add_blank_slide(prs)\n"
            "add_slide_title(slide, 'My title')\n"
            "add_bullets(slide, LEFT_MARGIN, ...)\n"
            "prs.save('output.pptx')"
        ),
    )

    add_callout_box(
        slide,
        left=Inches(7.0), top=Inches(1.4),
        width=Inches(5.5), height=Inches(1.4),
        label="Takeaway: ",
        text="callout boxes highlight key messages with a clean white background and thin border.",
        font_size=18,
    )

    add_textbox(
        slide, Inches(7.0), Inches(3.3), Inches(5.5), Inches(1.2),
        "Both elements are available as\nhelper functions in pptx_helpers.py",
        font_size=BODY_SIZE - 2, color=GREY,
    )


def slide_06_two_column(prs):
    """Two-column layout example."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Side-by-side layouts compare results or show figure + text")
    add_accent_line(slide)

    # Left column — placeholder "figure"
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Inches(1.5), Inches(5.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    shape.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    shape.line.width = Pt(1.5)
    add_textbox(
        slide, LEFT_MARGIN + Inches(1.3), Inches(3.2),
        Inches(3.0), Inches(0.6),
        "[Figure / Chart]",
        font_size=SUBTITLE_SIZE, color=GREY, alignment=PP_ALIGN.CENTER,
    )

    # Right column — bullets
    add_bullets(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0), [
        ("Left panel: ", "figure or chart occupies ~half the slide"),
        ("Right panel: ", "interpretation bullets summarize the finding"),
        "This layout works well for results slides",
        ("Tip: ", "keep to 3-4 bullets beside a figure"),
    ])


def slide_07_conclusions(prs):
    """Conclusions slide — last main slide, stays visible during Q&A."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Conclusions")
    add_accent_line(slide)
    add_bullets(slide, LEFT_MARGIN, TOP_MARGIN, CONTENT_W, Inches(4.5), [
        ("1. ", "Action titles make every slide self-contained and scannable"),
        ("2. ", "Blue accent prefixes replace bold for clean visual hierarchy"),
        ("3. ", "The style is fully customizable via slide-style.yaml"),
        ("4. ", "End with Conclusions, not 'Thank You' — this slide stays up during Q&A"),
    ])


def slide_08_supplementary(prs):
    """Supplementary / backup slide."""
    slide = add_blank_slide(prs)
    make_content_slide(slide, "Supplementary: full style configuration reference")
    add_accent_line(slide)

    data = [
        ["Setting", "Default", "Override via"],
        ["Font", "Avenir Light", "typography.body_font"],
        ["Code font", "PT Mono", "typography.code_font"],
        ["Accent color", "#036DEA", "colors.accent"],
        ["Bold", "Never", "style.use_bold"],
        ["Bullet emphasis", "Blue prefix", "style.bullet_emphasis"],
        ["Callout bg", "White + border", "style.callout_bg"],
        ["Title slide bg", "White", "style.title_slide_bg"],
    ]
    add_table(
        slide,
        left=Inches(1.5), top=Inches(1.5),
        width=Inches(10.0), height=Inches(4.5),
        rows=8, cols=3,
        data=data,
    )


def main():
    prs = new_presentation()

    slide_01_title(prs)
    slide_02_content(prs)
    slide_03_figure_placeholder(prs)
    slide_04_table(prs)
    slide_05_code_and_callout(prs)
    slide_06_two_column(prs)
    slide_07_conclusions(prs)
    slide_08_supplementary(prs)

    out = Path(__file__).resolve().parent / "sample-deck.pptx"
    prs.save(str(out))
    print(f"Generated {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
