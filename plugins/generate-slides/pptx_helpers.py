"""Reusable PPTX generation helpers for Claude Code generate-slides skill.

Provides a clean, minimal slide style (16:9 widescreen, Avenir Light font,
white backgrounds with sparing blue accents). All functions operate on a
python-pptx Presentation object.

Style profile extracted from Y.J. Kim's exemplar decks (2024-2025):
- Font: Avenir Light (body/titles), PT Mono (code)
- Colors: Black text, single blue (#036DEA) accent, no colored boxes
- Bold: Never — hierarchy through font size and spacing only
- Title position: top-aligned (~0.15")
- Font sizes: title 40pt, subtitle 28pt, body 22pt, footnote 14pt

Usage from a generated script:

    from pptx_helpers import *

    prs = new_presentation()
    slide = add_blank_slide(prs)
    add_slide_title(slide, "My Title", subtitle="Optional subtitle")
    add_bullets(slide, LEFT_MARGIN, Inches(1.5), CONTENT_W, Inches(4.0), [
        ("Key term: ", "rest of the bullet"),
        "Simple bullet text",
    ])
    prs.save("output.pptx")
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# Color Palette — intentionally minimal
# ============================================================================
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x03, 0x6D, 0xEA)          # primary accent (#036DEA)
GREEN = RGBColor(0x05, 0x96, 0x69)          # positive values
RED = RGBColor(0xDC, 0x26, 0x26)            # negative values
GREY = RGBColor(0x6B, 0x72, 0x80)           # secondary text / footnotes
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)    # table header bg
TABLE_BORDER = RGBColor(0xD1, 0xD5, 0xDB)

# ============================================================================
# Layout Constants (16:9 widescreen, standard PPTX)
# ============================================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
LEFT_MARGIN = Inches(0.8)
TOP_MARGIN = Inches(1.2)           # content starts below title + accent
CONTENT_W = Inches(11.7)

# ============================================================================
# Typography
# ============================================================================
FONT_NAME = "Avenir Light"
CODE_FONT = "PT Mono"

# Font sizes (scaled for 13.333" wide canvas)
TITLE_SIZE = 40
SUBTITLE_SIZE = 28
BODY_SIZE = 22
FOOTNOTE_SIZE = 14
TABLE_SIZE = 16
CODE_SIZE = 12


# ============================================================================
# Presentation & Slide Setup
# ============================================================================

def new_presentation():
    """Create a new 16:9 widescreen presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================================
# Text Primitives
# ============================================================================

def add_textbox(slide, left, top, width, height, text, font_size=BODY_SIZE,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, line_spacing=1.15):
    """Add a simple single-run text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_textbox(slide, left, top, width, height, paragraphs_data,
                     font_name=FONT_NAME, line_spacing=1.15):
    """Add a text box with multiple paragraphs, each with multiple styled runs.

    paragraphs_data: list of dicts, each with:
        - runs: list of dicts {text, font_size, bold, color, italic}
        - alignment: PP_ALIGN (default LEFT)
        - space_after: Pt value (default 4)
        - bullet: bool (prepend bullet character)
        - indent: bool (indent the paragraph)
        - level: int (indentation level)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, pdata in enumerate(paragraphs_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pdata.get("alignment", PP_ALIGN.LEFT)
        p.space_after = pdata.get("space_after", Pt(4))

        if pdata.get("indent"):
            p.level = pdata.get("level", 1)

        # Prepend bullet character to first run if requested
        runs_data = list(pdata.get("runs", []))
        if pdata.get("bullet") and runs_data:
            runs_data[0] = dict(runs_data[0])
            runs_data[0]["text"] = "\u2022  " + runs_data[0]["text"]

        for rdata in runs_data:
            run = p.add_run()
            run.text = rdata.get("text", "")
            run.font.size = Pt(rdata.get("font_size", BODY_SIZE))
            run.font.bold = rdata.get("bold", False)
            run.font.italic = rdata.get("italic", False)
            run.font.color.rgb = rdata.get("color", BLACK)
            run.font.name = font_name

        if line_spacing != 1.0:
            fs = runs_data[0].get("font_size", BODY_SIZE) if runs_data else BODY_SIZE
            p.line_spacing = Pt(fs * line_spacing)

    return txBox


# ============================================================================
# Slide Components
# ============================================================================

def add_slide_title(slide, title, subtitle=None, color=BLACK):
    """Add a title (and optional subtitle) to a content slide.

    Title sits flush near the top (~0.15") for maximum content area.
    """
    add_textbox(slide, LEFT_MARGIN, Inches(0.15), CONTENT_W, Inches(0.65),
                title, font_size=TITLE_SIZE, bold=False, color=color)
    if subtitle:
        add_textbox(slide, LEFT_MARGIN, Inches(0.75), CONTENT_W, Inches(0.4),
                    subtitle, font_size=SUBTITLE_SIZE, color=GREY, bold=False)


def add_accent_line(slide, left=None, top=Inches(1.05), width=Inches(4.0),
                    color=BLUE):
    """Add a thin horizontal accent line (default: blue, below title)."""
    left = left or LEFT_MARGIN
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_figure(slide, fig_path, left, top, width=None, height=None):
    """Add an image. Returns the picture shape, or None if file missing."""
    fig_path = Path(fig_path)
    if not fig_path.exists():
        print(f"WARNING: Figure not found: {fig_path}")
        return None
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(fig_path), left, top, **kwargs)


def add_code_box(slide, left, top, width, height, text, font_size=CODE_SIZE):
    """Add a monospace text box with light grey background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = CODE_FONT
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = Pt(font_size * 1.3)

    return shape


def add_callout_box(slide, left, top, width, height, label, text,
                    label_color=BLUE, text_color=BLACK, font_size=20):
    """Add a callout/takeaway box with subtle styling.

    Uses a thin left border (blue accent line) instead of a colored
    background, keeping the slide clean and readable.
    """
    # White box with thin grey border
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.color.rgb = label_color
    run.font.name = FONT_NAME
    run.font.bold = False
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.name = FONT_NAME
    run.font.bold = False
    return shape


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, header_color=LIGHT_GREY, font_size=TABLE_SIZE,
              bold_rows=None, highlight_cells=None):
    """Add a styled table.

    Args:
        data: list of lists (row-major), first row is header.
        col_widths: list of Inches/Emu per column.
        bold_rows: set of row indices to render in medium weight (not true bold).
        highlight_cells: dict of (row, col) -> RGBColor for text color.
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    bold_rows = bold_rows or set()
    highlight_cells = highlight_cells or {}

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_NAME
                paragraph.font.bold = False
                if (r, c) in highlight_cells:
                    paragraph.font.color.rgb = highlight_cells[(r, c)]
                elif r == 0:
                    paragraph.font.color.rgb = BLACK
                else:
                    paragraph.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ============================================================================
# Bullet Helpers
# ============================================================================

def make_bullet_paragraphs(items, font_size=BODY_SIZE, color=BLACK):
    """Convert a list of items into paragraphs_data for add_rich_textbox.

    Items can be:
        - str: simple bullet
        - (emphasis_part, rest): tuple with emphasis prefix (rendered in blue,
          not bold — consistent with the no-bold style)
    """
    paragraphs = []
    for item in items:
        if isinstance(item, str):
            paragraphs.append({
                "runs": [{"text": item, "font_size": font_size, "color": color}],
                "bullet": True,
                "space_after": Pt(8),
            })
        elif isinstance(item, tuple) and len(item) == 2:
            paragraphs.append({
                "runs": [
                    {"text": item[0], "font_size": font_size, "color": BLUE},
                    {"text": item[1], "font_size": font_size, "color": color},
                ],
                "bullet": True,
                "space_after": Pt(8),
            })
        elif isinstance(item, dict):
            paragraphs.append(item)
    return paragraphs


def add_bullets(slide, left, top, width, height, items, font_size=BODY_SIZE,
                color=BLACK):
    """Convenience: add a bulleted list to a slide.

    Items can be str or (emphasis_prefix, rest) tuples.
    Emphasis prefix is rendered in blue accent color (no bold).
    """
    paras = make_bullet_paragraphs(items, font_size=font_size, color=color)
    return add_rich_textbox(slide, left, top, width, height, paras)


# ============================================================================
# Composite Slide Templates
# ============================================================================

def make_title_slide(slide, title, author, affiliation, date):
    """Build a white-background title slide.

    Clean layout: title at top, accent line, then author/affiliation/date.
    """
    set_slide_bg(slide, WHITE)
    add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.6),
                title, font_size=TITLE_SIZE, bold=False, color=BLACK,
                line_spacing=1.25)
    add_accent_line(slide, Inches(1.0), Inches(4.3), Inches(2.0))
    add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.0), Inches(0.4),
                author, font_size=SUBTITLE_SIZE, color=BLACK)
    add_textbox(slide, Inches(1.0), Inches(5.15), Inches(11.0), Inches(0.4),
                affiliation, font_size=FOOTNOTE_SIZE + 4, color=GREY)
    add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11.0), Inches(0.4),
                date, font_size=FOOTNOTE_SIZE + 4, color=GREY)


def make_content_slide(slide, title, subtitle=None):
    """Set up a standard content slide with title + accent line."""
    add_slide_title(slide, title, subtitle=subtitle)
    add_accent_line(slide)


def make_section_slide(slide, section_title, bg_color=None):
    """Build a section divider slide (large centered text)."""
    if bg_color:
        set_slide_bg(slide, bg_color)
    text_color = WHITE if bg_color else BLACK
    add_textbox(slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0),
                section_title, font_size=TITLE_SIZE + 4, bold=False,
                color=text_color, alignment=PP_ALIGN.CENTER)
